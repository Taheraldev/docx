import io
from docx import Document
from docx.shared import Pt
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Pt as pptxPt
from googletrans import Translator
import arabic_reshaper
from bidi.algorithm import get_display
from telegram import Update
from telegram.ext import Updater, CommandHandler, MessageHandler, Filters, CallbackContext

# أدخل توكن البوت الخاص بك هنا
TOKEN = '5284087690:AAGwKfPojQ3c-SjCHSIdeog-yN3-4Gpim1Y'

# تهيئة مترجم جوجل
translator = Translator()

# خيار إعادة تشكيل النصوص العربية
apply_arabic_processing = False

# تحديد الخط العربي الافتراضي
ARABIC_FONT = "Times New Roman"

def process_arabic(text: str) -> str:
    """
    إعادة تشكيل النص العربي إذا لزم الأمر.
    """
    if apply_arabic_processing:
        reshaped_text = arabic_reshaper.reshape(text)
        bidi_text = get_display(reshaped_text)
        return bidi_text
    else:
        return text

def set_paragraph_rtl(paragraph):
    """
    ضبط اتجاه الكتابة من اليمين لليسار باستخدام خصائص XML.
    """
    pPr = paragraph._p.get_or_add_pPr()
    bidi = OxmlElement('w:bidi')
    bidi.set(qn('w:val'), "1")
    pPr.append(bidi)

def translate_paragraph(paragraph):
    """
    ترجمة النص داخل كل run من الفقرة دون إزالة العناصر غير النصية مثل الصور.
    """
    for run in paragraph.runs:
        if run.text.strip():
            translated_text = translator.translate(run.text, src='en', dest='ar').text
            translated_text = process_arabic(translated_text)
            run.text = translated_text
            # تعيين الخط العربي وحجم الخط للـ run
            run.font.name = ARABIC_FONT
            run.font.size = Pt(14)
    # ضبط اتجاه الفقرة إلى RTL
    set_paragraph_rtl(paragraph)

def translate_docx(file_bytes: bytes) -> io.BytesIO:
    """
    ترجمة محتوى DOCX بما في ذلك الفقرات والجداول دون فقدان الصور.
    """
    document = Document(io.BytesIO(file_bytes))
    
    # ترجمة الفقرات الرئيسية
    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            translate_paragraph(paragraph)
    
    # ترجمة النصوص داخل الجداول
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    if paragraph.text.strip():
                        translate_paragraph(paragraph)
    
    output = io.BytesIO()
    document.save(output)
    output.seek(0)
    return output

def translate_pptx(file_bytes: bytes) -> io.BytesIO:
    """
    ترجمة محتوى PPTX بما في ذلك النصوص داخل مربعات النص والجداول.
    """
    prs = Presentation(io.BytesIO(file_bytes))
    
    for slide in prs.slides:
        for shape in slide.shapes:
            # ترجمة المربعات النصية
            if hasattr(shape, "text") and shape.text.strip():
                translated_text = translator.translate(shape.text, src='en', dest='ar').text
                translated_text = process_arabic(translated_text)
                if shape.has_text_frame:
                    shape.text = translated_text
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = ARABIC_FONT
                            run.font.size = pptxPt(24)
            # ترجمة الجداول في PPTX
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            translated_text = translator.translate(cell.text, src='en', dest='ar').text
                            translated_text = process_arabic(translated_text)
                            cell.text = translated_text
                            if cell.text_frame:
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.name = ARABIC_FONT
                                        run.font.size = pptxPt(18)
    
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

def start(update: Update, context: CallbackContext) -> None:
    update.message.reply_text("مرحباً! أرسل لي ملف DOCX أو PPTX للترجمة من الإنجليزية إلى العربية.")

def handle_file(update: Update, context: CallbackContext) -> None:
    document_file = update.message.document
    file_name = document_file.file_name.lower()
    
    file = document_file.get_file()
    file_bytes = file.download_as_bytearray()

    if file_name.endswith('.docx'):
        translated_file = translate_docx(file_bytes)
        update.message.reply_document(document=translated_file, filename="translated.docx")
    elif file_name.endswith('.pptx'):
        translated_file = translate_pptx(file_bytes)
        update.message.reply_document(document=translated_file, filename="translated.pptx")
    else:
        update.message.reply_text("صيغة الملف غير مدعومة. الرجاء إرسال ملف DOCX أو PPTX.")

def main():
    updater = Updater(TOKEN, use_context=True)
    dp = updater.dispatcher
    
    dp.add_handler(CommandHandler("start", start))
    dp.add_handler(MessageHandler(Filters.document, handle_file))
    
    updater.start_polling()
    updater.idle()

if __name__ == '__main__':
    main()
